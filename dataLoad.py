from pptx import Presentation
import threading
import pymupdf
import os
import time



base_dir = os.getcwd()
class dataLoad():
    def getModulePath(self, mod):
        if mod=="FIT152":
            return "data/modules/FIT152/"
        elif mod=="ISP152":
            return "data/modules/ISP152"
        elif mod=="IDB152":
            return "data/modules/IDB152"
        elif mod=="OOP152":
            return "data/modules/OOP152"
        elif mod=="SEN152":
            return "data/modules/SEN152"
        elif mod=="TAS152":
            return "data/modules/TAS152"
    
    def getPresText(self, file_path):
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return "\n".join(text)

    def getPDFtxt(self, file_path):
        doc = pymupdf.open(file_path)
        allText = ""
        for page in doc:
            allText += page.get_text()
        return allText
    
    def loadPowerP(self, modPath):
        fold = os.path.join(modPath , "ppSlides")
        ls = os.listdir(fold) 
        ppText = [] 
        for file in ls:
            if not file.lower().endswith('.pptx'):
                continue
            fp = os.path.join(base_dir,fold, file)
            txt = self.getPresText(fp)
            ppText.append({"File": file, "Content": txt})
        return ppText
    
    def loadPDF(self, modPath):
        fold = os.path.join(modPath , "studyGuides")
        ls = os.listdir(fold)
        pdftxt = []
        for file in ls:
            if not file.lower().endswith('.pdf'):
                continue
            fp = os.path.join(base_dir,fold, file)
            txt = self.getPDFtxt(fp)
            pdftxt.append({"File": file, "Content": txt})
        return pdftxt

    def loadAll(self, module):
        mod = self.getModulePath(module)
        start_time = time.time()
        results = [None, None]
        
        threads = [
            threading.Thread(
                target=lambda: results.__setitem__(0, self.loadPowerP(mod))
            ),
            threading.Thread(
                target=lambda: results.__setitem__(1, self.loadPDF(mod))
            )
        ]
        
        # Start all threads
        for t in threads:
            t.start()

        # Wait for all to finish
        for t in threads:
            t.join()
            
        ppt = results[0]
        pdf = results[1]
        
        
        sum_val = 0
        for i in range(1000000):
            sum_val += i

        end_time = time.time()
        elapsed_time = end_time - start_time

        print(f"-----Loading All data time: {round(elapsed_time,3)} seconds...-----")
        return ppt, pdf
    
    


